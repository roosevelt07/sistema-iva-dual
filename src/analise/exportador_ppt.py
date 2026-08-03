"""Exportação do relatório de análise em PowerPoint (.pptx)."""

import io
from decimal import Decimal
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Cm, Pt
    from lxml import etree

    PPTX_DISPONIVEL = True
except ImportError:
    PPTX_DISPONIVEL = False

from src.analise.decisao import ResultadoAnalise
from src.analise.formatters import fmt_moeda, fmt_percentual

LOGO_PATH = Path(__file__).parent.parent.parent / "assets" / "logo.png"

_VERDE = RGBColor(0x00, 0xB2, 0xA9) if PPTX_DISPONIVEL else None  # teal Guerra
_VERDE_DK = RGBColor(0x00, 0x7F, 0x78) if PPTX_DISPONIVEL else None  # teal escuro
_PRETO = RGBColor(0x2C, 0x2C, 0x2C) if PPTX_DISPONIVEL else None
_CINZA = RGBColor(0x6B, 0x72, 0x80) if PPTX_DISPONIVEL else None
_BRANCO = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_DISPONIVEL else None
_LIGHT = RGBColor(0xF4, 0xFA, 0xFA) if PPTX_DISPONIVEL else None

_LARGURA = Cm(33.87) if PPTX_DISPONIVEL else None
_ALTURA = Cm(19.05) if PPTX_DISPONIVEL else None


def _fundo(slide, cor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = cor


def _caixa_texto(slide, texto, left, top, width, height, tamanho=18, cor=None,
                  negrito=False, italico=False, alinhamento=PP_ALIGN.LEFT, fonte="Calibri"):
    caixa = slide.shapes.add_textbox(left, top, width, height)
    tf = caixa.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alinhamento
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.name = fonte
    run.font.bold = negrito
    run.font.italic = italico
    if cor is not None:
        run.font.color.rgb = cor
    return caixa


def _retangulo(slide, left, top, width, height, cor):
    forma = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    forma.fill.solid()
    forma.fill.fore_color.rgb = cor
    forma.line.fill.background()
    forma.shadow.inherit = False
    return forma


def _adicionar_sombra(forma) -> None:
    """Sombra suave via XML — não exposta pela API alta do python-pptx.

    Cosmético: se algo falhar (versão de biblioteca diferente etc.), a
    exportação segue sem sombra em vez de quebrar.
    """
    try:
        forma.shadow.inherit = False
        effect_lst = forma._element.spPr.find(qn("a:effectLst"))
        outer_shdw = etree.SubElement(effect_lst, qn("a:outerShdw"))
        outer_shdw.set("blurRad", "90000")
        outer_shdw.set("dist", "30000")
        outer_shdw.set("dir", "5400000")
        outer_shdw.set("rotWithShape", "0")
        srgb_clr = etree.SubElement(outer_shdw, qn("a:srgbClr"))
        srgb_clr.set("val", "2C2C2C")
        alpha = etree.SubElement(srgb_clr, qn("a:alpha"))
        alpha.set("val", "35000")
    except Exception:
        pass


def _cartao(slide, left, top, width, height, valor_texto, label_texto):
    cartao = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    cartao.fill.solid()
    cartao.fill.fore_color.rgb = _BRANCO
    cartao.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    cartao.line.width = Pt(0.75)
    _adicionar_sombra(cartao)

    tf = cartao.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.3)
    tf.margin_right = Cm(0.3)

    p_valor = tf.paragraphs[0]
    p_valor.alignment = PP_ALIGN.CENTER
    run_valor = p_valor.add_run()
    run_valor.text = valor_texto
    run_valor.font.size = Pt(24)
    run_valor.font.bold = True
    run_valor.font.color.rgb = _VERDE_DK

    p_label = tf.add_paragraph()
    p_label.alignment = PP_ALIGN.CENTER
    run_label = p_label.add_run()
    run_label.text = label_texto
    run_label.font.size = Pt(13)
    run_label.font.color.rgb = _CINZA


def _cabecalho_faixa(slide, titulo, cor_fundo, cor_texto):
    _retangulo(slide, 0, 0, _LARGURA, Cm(2.2), cor_fundo)
    _caixa_texto(
        slide, titulo, Cm(1.5), Cm(0.5), Cm(28), Cm(1.4),
        tamanho=26, cor=cor_texto, negrito=True,
    )


def _adicionar_logo(slide, left, top, altura):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, height=altura)


def _slide_capa(prs, layout, resultado: ResultadoAnalise, nome_cliente, data_relatorio):
    slide = prs.slides.add_slide(layout)
    _fundo(slide, _PRETO)

    # Faixa teal lateral esquerda — ocupa toda a altura do slide.
    _retangulo(slide, 0, 0, Cm(1.2), _ALTURA, _VERDE)

    _adicionar_logo(slide, Cm(2), Cm(1), Cm(3))

    _caixa_texto(
        slide, "Análise de Impacto", Cm(2.5), Cm(6), Cm(28), Cm(3),
        tamanho=32, cor=_BRANCO, negrito=True,
    )
    _caixa_texto(
        slide, "Reforma Tributária — LC 214/2025", Cm(2.5), Cm(9.2), Cm(28), Cm(1.5),
        tamanho=20, cor=_VERDE,
    )
    _retangulo(slide, Cm(2.5), Cm(11), Cm(15), Pt(1.5), _VERDE)

    _caixa_texto(
        slide, f"{nome_cliente or '—'}", Cm(2.5), Cm(11.4), Cm(20), Cm(1),
        tamanho=15, cor=_BRANCO,
    )
    _caixa_texto(
        slide, f"{data_relatorio}", Cm(2.5), Cm(12.2), Cm(20), Cm(0.8),
        tamanho=11, cor=_CINZA,
    )

    _caixa_texto(
        slide, "Simulação baseada na LC 214/2025 — valores estimados.",
        Cm(2.5), Cm(17.5), Cm(28), Cm(1),
        tamanho=9, cor=_CINZA, italico=True,
    )
    return slide


def _slide_resumo_executivo(prs, layout, resultado: ResultadoAnalise):
    cenario_a = resultado.cenarios.cenario_a
    cenario_b = resultado.cenarios.cenario_b

    slide = prs.slides.add_slide(layout)
    _fundo(slide, _LIGHT)
    _cabecalho_faixa(slide, "Resumo Executivo", _VERDE, _BRANCO)

    cartoes = [
        (fmt_moeda(cenario_a.carga_liquida), "Carga Atual (líquida)"),
        (fmt_moeda(cenario_b.carga_liquida), "Carga IVA (líquida)"),
        (fmt_moeda(abs(resultado.delta_absoluto)), "Economia"),
    ]
    largura_cartao = Cm(9.5)
    espaco = Cm(1.2)
    inicio_x = Cm(1.9)
    for i, (valor, label) in enumerate(cartoes):
        left = inicio_x + i * (largura_cartao + espaco)
        _cartao(slide, left, Cm(4), largura_cartao, Cm(6), valor, label)

    sinal = "economia" if resultado.delta_absoluto > 0 else "aumento"
    _caixa_texto(
        slide,
        f"A migração para o regime regular do IBS/CBS representa {sinal} de "
        f"{fmt_percentual(abs(resultado.delta_percentual))} em relação à carga atual.",
        Cm(1.9), Cm(12), Cm(30), Cm(2),
        tamanho=16, cor=_PRETO,
    )
    return slide


def _slide_comparativo_carga(prs, layout, resultado: ResultadoAnalise):
    cenario_a = resultado.cenarios.cenario_a
    cenario_b = resultado.cenarios.cenario_b

    slide = prs.slides.add_slide(layout)
    _fundo(slide, _BRANCO)
    _cabecalho_faixa(slide, "Comparativo de Carga", _PRETO, _BRANCO)

    linhas_dados = [
        ("Carga bruta", cenario_a.carga_bruta, cenario_b.carga_bruta),
        ("Crédito aproveitado", cenario_a.credito_aproveitado, cenario_b.credito_aproveitado),
        ("Carga líquida", cenario_a.carga_liquida, cenario_b.carga_liquida),
    ]

    # Gráfico de barras nativo (editável no PowerPoint, não uma imagem).
    chart_data = CategoryChartData()
    chart_data.categories = [nome for nome, _, _ in linhas_dados]
    chart_data.add_series("Atual", tuple(float(valor_a) for _, valor_a, _ in linhas_dados))
    chart_data.add_series("IVA", tuple(float(valor_b) for _, _, valor_b in linhas_dados))
    grafico_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Cm(1.5), Cm(3), Cm(19), Cm(13), chart_data
    )
    chart = grafico_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = _PRETO
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = _VERDE

    # Tabela lateral.
    tabela_shape = slide.shapes.add_table(len(linhas_dados) + 1, 4, Cm(21.5), Cm(3), Cm(11), Cm(8))
    tabela = tabela_shape.table
    for c, texto in enumerate(["Item", "Atual", "IVA", "Economia"]):
        celula = tabela.cell(0, c)
        celula.text = texto
        celula.fill.solid()
        celula.fill.fore_color.rgb = _VERDE
        for run in celula.text_frame.paragraphs[0].runs:
            run.font.color.rgb = _BRANCO
            run.font.bold = True
            run.font.size = Pt(11)
    for i, (nome, valor_a, valor_b) in enumerate(linhas_dados, start=1):
        valores_linha = [nome, fmt_moeda(valor_a), fmt_moeda(valor_b), fmt_moeda(valor_a - valor_b)]
        for c, texto in enumerate(valores_linha):
            celula = tabela.cell(i, c)
            celula.text = texto
            celula.fill.solid()
            celula.fill.fore_color.rgb = _LIGHT if i % 2 == 0 else _BRANCO
            for run in celula.text_frame.paragraphs[0].runs:
                run.font.size = Pt(10)
    return slide


def _slide_atividades(prs, layout, resultado: ResultadoAnalise):
    slide = prs.slides.add_slide(layout)
    _fundo(slide, _LIGHT)
    _cabecalho_faixa(slide, "Atividades Declaradas", _VERDE, _BRANCO)

    # ResultadoAnalise/ContextoCalculo não carregam dados por atividade
    # hoje (isso vive em ExtratoPGDAS/DadosNormalizados, camada do
    # parser). Guard defensivo: funciona automaticamente se uma versão
    # futura passar a anexar esses dados.
    blocos_atividade = getattr(resultado.ctx, "blocos_atividade", None)
    if not blocos_atividade:
        _caixa_texto(
            slide, "Dados disponíveis no Modo 3 — Extrato PGDAS-D.",
            Cm(1.9), Cm(8), Cm(28), Cm(2),
            tamanho=18, cor=_PRETO,
        )
        return slide

    tabela_shape = slide.shapes.add_table(
        len(blocos_atividade) + 1, 5, Cm(1.5), Cm(3), Cm(30), Cm(2 + len(blocos_atividade))
    )
    tabela = tabela_shape.table
    for c, texto in enumerate(["Atividade", "Anexo", "Receita", "DAS", "Obs."]):
        celula = tabela.cell(0, c)
        celula.text = texto
        celula.fill.solid()
        celula.fill.fore_color.rgb = _VERDE
        for run in celula.text_frame.paragraphs[0].runs:
            run.font.color.rgb = _BRANCO
            run.font.bold = True
    for i, bloco in enumerate(blocos_atividade, start=1):
        valores_linha = [
            getattr(bloco, "nome", ""),
            str(getattr(bloco, "anexo", "")),
            fmt_moeda(getattr(bloco, "receita_atividade", Decimal("0"))),
            fmt_moeda(getattr(bloco, "das", Decimal("0"))),
            getattr(bloco, "observacao", ""),
        ]
        for c, texto in enumerate(valores_linha):
            celula = tabela.cell(i, c)
            celula.text = texto
            celula.fill.solid()
            celula.fill.fore_color.rgb = _LIGHT if i % 2 == 0 else _BRANCO

    anexo_predominante = getattr(resultado, "anexo_predominante", None)
    if anexo_predominante:
        _caixa_texto(
            slide, f"Anexo predominante: {anexo_predominante}",
            Cm(1.5), Cm(17), Cm(20), Cm(1),
            tamanho=12, cor=_CINZA,
        )
    return slide


def _slide_recomendacao(prs, layout, resultado: ResultadoAnalise):
    slide = prs.slides.add_slide(layout)
    migrar = resultado.recomendacao == "MIGRAR_REGIME_REGULAR"
    _fundo(slide, _VERDE if migrar else _PRETO)

    _adicionar_logo(slide, Cm(1.5), Cm(1), Cm(1.8))

    texto_recomendacao = "✓ MIGRAR PARA REGIME REGULAR" if migrar else "→ MANTER REGIME ATUAL"
    _caixa_texto(
        slide, texto_recomendacao, Cm(2), Cm(6), Cm(30), Cm(3),
        tamanho=36, cor=_BRANCO, negrito=True, alinhamento=PP_ALIGN.CENTER,
    )
    subtexto = f"Delta: {fmt_moeda(resultado.delta_absoluto)} ({fmt_percentual(resultado.delta_percentual)})"
    _caixa_texto(
        slide, subtexto, Cm(2), Cm(9.3), Cm(30), Cm(2),
        tamanho=20, cor=_BRANCO, alinhamento=PP_ALIGN.CENTER,
    )

    if resultado.alertas:
        caixa = slide.shapes.add_textbox(Cm(4), Cm(12), Cm(26), Cm(5))
        tf = caixa.text_frame
        tf.word_wrap = True
        for i, alerta in enumerate(resultado.alertas):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {alerta}"
            p.font.size = Pt(12)
            p.font.color.rgb = _BRANCO

    _caixa_texto(
        slide, "Simulação baseada na LC 214/2025 — valores estimados.",
        Cm(2), Cm(17.8), Cm(28), Cm(1),
        tamanho=9, cor=_BRANCO, italico=True,
    )
    return slide


def gerar_ppt(resultado: ResultadoAnalise, nome_cliente: str, data_relatorio) -> bytes:
    """Gera a apresentação PPT (5 slides) e retorna os bytes do arquivo .pptx."""
    if not PPTX_DISPONIVEL:
        raise RuntimeError("python-pptx não está instalado.")

    prs = Presentation()
    prs.slide_width = _LARGURA
    prs.slide_height = _ALTURA
    layout_branco = prs.slide_layouts[6]

    _slide_capa(prs, layout_branco, resultado, nome_cliente, data_relatorio)
    _slide_resumo_executivo(prs, layout_branco, resultado)
    _slide_comparativo_carga(prs, layout_branco, resultado)
    _slide_atividades(prs, layout_branco, resultado)
    _slide_recomendacao(prs, layout_branco, resultado)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
